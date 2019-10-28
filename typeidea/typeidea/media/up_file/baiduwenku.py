import requests
from selenium import webdriver
import time
import os
from pptx import Presentation  # 生成ppt需要的包


def get_html(url, return_type="text"):  # 获取网页页面，返回文本文件或者二进制文件
    r = requests.get(url)
    r.raise_for_status()
    r.encoding = r.apparent_encoding
    if return_type == "text":
        return r.text
    else:
        return r.content


def get_picture_url(url):  # 得到该页面的每页图片的url
    opt = webdriver.FirefoxOptions()  # 设置驱动的浏览器是火狐
    opt.set_headless()
    driver = webdriver.Chrome(options=opt)  # 启动火狐，无界面模式
    driver.get(url)  # 打开url
    time.sleep(3)  # 休息3秒，等待加载
    # xpath寻找图片所在的tag
    img_tags = driver.find_elements_by_xpath("//div[@class='ppt-image-wrap']/img")
    img_urls = []
    for img_tag in img_tags:  # 提取url
        if img_tag.get_attribute("src"):
            img_urls.append(img_tag.get_attribute("src"))
        else:
            img_urls.append(img_tag.get_attribute("data-src"))
    return img_urls  # 返回


def download_pictures(url):  # 下载图片，并存储
    img_urls = get_picture_url(url)  # 得到图片url
    if not os.path.exists("./pictures"):  # 判断是否有存储路径
        os.makedirs("./pictures")
    path_save = []  # 储存图片的名字和路径，在合成ppt时保持正确的顺序
    for i in range(len(img_urls)):
        img = get_html(img_urls[i], return_type="img")  # 得到图片的二进制文件
        with open("./pictures/num_%s.jpg" % str(i), "wb") as f:  # 写出
            path_save.append("./pictures/num_%s.jpg" % str(i))
            f.write(img)
    return path_save  # 返回正确顺序的图片路径和名字


def download_word(url):  # 如果是word， 则调用这个函数
    opt = webdriver.ChromeOptions()
    opt.set_headless()
    driver = webdriver.Chrome(options=opt)
    driver.get(url)
    time.sleep(3)
    # 找到word中文字段落所在的标签
    p_tags = driver.find_elements_by_xpath("//div[@class='reader-txt-layer']/div/p")
    all_text = ""
    for p in p_tags:  # 循环所有的tag
        text = p.text  # 得到文本
        if text == "":  # 如果是空，则表示换行（观察html得出的结论）
            all_text = all_text + "\n"
        else:
            all_text = all_text + text
    return all_text  # 返回拼接好的字符串


def generate_ppt(picture_paths):  # 如果是ppt， 在上面下载好图片后，这个函数嫁给你图片重新合成ppt
    prs = Presentation('test.pptx')  # 这个是一个ppt模板，里面只有一个图片的占位
    for picture_path in picture_paths:
        # 按照第一个模板创建 一张幻灯片
        oneSlide = prs.slides.add_slide(prs.slide_layouts[0])
        # 获取模板可填充的所有位置
        body_shapes = oneSlide.shapes.placeholders
        for index, body_shape in enumerate(body_shapes):
            if index == 0:
                body_shape.insert_picture(picture_path)

    # 对ppt的修改
    prs.save('baiduwenku.pptx')
    # 这个函数具体的不是很懂，在网上copy的


def write_to_file(url, type):  # 写入函数
    if type == "ppt":
        path_history = download_pictures(url)
        generate_ppt(path_history)
        print("download successfully!")
    elif type == "word":
        text = download_word(url)
        with open("aa.doc", "w", encoding="utf-8") as f:
            f.write(text)
        print("download successfully!")
    else:
        print("Does not support this type of file now.")


if __name__ == "__main__":  # 主函数
    url = "https://wenku.baidu.com/view/926285d680eb6294dd886c73.html?from=search"
    return_type = "word"
    write_to_file(url, return_type)
    # 1， 判断是ppt还是word，
    # 2， 两种不同的文件调用不同的函数


